from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from time import *

begin_time = time()
MdList = []
PostList = []
with open("ppp.txt", "r") as f:
    for line in f.readlines():
        line = line.strip('\n')
        if '/' in line:
            MdList.append(line)
        else:
            PostList.append(line)

TitleList = []
PlatformList = []
with open("WORDING.txt", "r") as f:
    for line in f.readlines():
        line = line.strip('\n')
        if 'Title:' in line:
            TitleList.append(line.strip('Title:'))
        else:
            PlatformList.append(line.strip('PlatformList:'))

with open("COMMENT.txt", "r") as f:
    Comment = f.read()

with open("URL.txt", "r") as f:
    URL = f.read()

# ADD A NEW SLIDE
prs = Presentation('final W12-W13 Weibo Skincare Mockup-DIOR Prestige Spring Rituals_0316.pptx')
slide = prs.slides.add_slide(prs.slide_layouts[11])

# SKINCARE WEIBO
SKINCAREWEIBO = slide.shapes.add_textbox(Cm(6.33), Cm(1.78), Cm(15.07), Cm(4.34))
SKINCAREWEIBO_Text = SKINCAREWEIBO.text_frame.add_paragraph()
SKINCAREWEIBO_Text.text = PlatformList[0]
SKINCAREWEIBO_Text.font.size = Pt(45)
SKINCAREWEIBO_Text.font.color.rgb = RGBColor(0, 0, 0)
SKINCAREWEIBO_Text.font.name = 'Century Gothic'
SKINCAREWEIBO_Text.alignment = PP_ALIGN.LEFT

# WEIBOPIC
img_path = 'PIC/WEIBO.png'
WEIBOPIC = slide.shapes.add_picture(img_path, Cm(2.64), Cm(3.26), Cm(2.98), Cm(2.39))

# MOCKUPPIC
img_path = 'PIC/L1T1.jpg'
MOCKUPPIC = slide.shapes.add_picture(img_path, Cm(4), Cm(10.4), Cm(11.48), Cm(6))

# Original Visual
img_path = 'PIC/L1T2-1.jpg'
OV1 = slide.shapes.add_picture(img_path, Cm(5.47), Cm(29.51), Cm(4.63), Cm(6.94))

# Original Visual
img_path = 'PIC/L1T2-2.jpg'
OV2 = slide.shapes.add_picture(img_path, Cm(10.52), Cm(29.51), Cm(4.63), Cm(6.94))

# P1
img_path = 'PIC/P1.gif'
P1 = slide.shapes.add_picture(img_path, Cm(41.33), Cm(14.29), Cm(2.5), Cm(4))
# P2
img_path = 'PIC/P2.jpg'
P2 = slide.shapes.add_picture(img_path, Cm(51.45), Cm(14.29), Cm(2.5), Cm(4))
# P3
img_path = 'PIC/P3.gif'
P3 = slide.shapes.add_picture(img_path, Cm(61.56), Cm(14.29), Cm(2.5), Cm(4))
# P4
img_path = 'PIC/P4.jpg'
P4 = slide.shapes.add_picture(img_path, Cm(41.33), Cm(20.95), Cm(2.5), Cm(4))
# P5
img_path = 'PIC/P5.jpg'
P5 = slide.shapes.add_picture(img_path, Cm(51.45), Cm(20.95), Cm(2.5), Cm(4))
# P6
img_path = 'PIC/P6.jpg'
P6 = slide.shapes.add_picture(img_path, Cm(61.56), Cm(20.95), Cm(2.5), Cm(4))
# P7
img_path = 'PIC/P7.jpg'
P7 = slide.shapes.add_picture(img_path, Cm(41.33), Cm(27.3), Cm(2.5), Cm(4))
# P8
img_path = 'PIC/P8.jpg'
P8 = slide.shapes.add_picture(img_path, Cm(51.45), Cm(27.3), Cm(2.5), Cm(4))
# P9
img_path = 'PIC/P9.jpg'
P9 = slide.shapes.add_picture(img_path, Cm(61.56), Cm(27.3), Cm(2.5), Cm(4))
# P10
img_path = 'PIC/L1T1.jpg'
P10 = slide.shapes.add_picture(img_path, Cm(40.37), Cm(33.6), Cm(6), Cm(3))

# Title
body_shape = slide.shapes.placeholders
Title = body_shape[0]
Title.left = Cm(27.6)
Title.top = Cm(0.8)
Title.height = Cm(3.23)
Title.width = Cm(44.61)
body_shape[0].text_frame.clear()
Title_Text = body_shape[0].text_frame.add_paragraph()
Title_Text.text = TitleList[0]
Title_Text.font.bold = False
Title_Text.font.italic = False
Title_Text.font.size = Pt(45)
Title_Text.font.name = 'Century Gothic'
Title_Text.font.underline = False

# RECTANGLE1 is the red Rectangle on the top left corner
RECTANGLE1 = slide.shapes
left = Cm(2.7)
top = Cm(8.9)
width = Cm(13.69)
height = Cm(8.9)
shape = RECTANGLE1.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, width, height
)
fill = shape.fill
fill.solid()
shape.fill.background()

line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
line.width = Pt(2.0)

# Mockup
Mockup = slide.shapes
left = Cm(2.7)
top = Cm(7.2)
width = Cm(13.69)
height = Cm(1.24)
shape = Mockup.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(166, 166, 166)

line = shape.line
line.color.rgb = RGBColor(166, 166, 166)
line.width = Pt(2.0)

# MockupTitle_Text
MockupTitle = slide.shapes.add_textbox(Cm(2.7), Cm(5.11), Cm(13.69), Cm(1.24))
MockupTitle_Text = MockupTitle.text_frame.add_paragraph()
MockupTitle_Text.text = 'MOCKUP'
MockupTitle_Text.font.bold = False
MockupTitle_Text.font.italic = False
MockupTitle_Text.font.size = Pt(35.4)
MockupTitle_Text.font.color.rgb = RGBColor(255, 255, 255)
MockupTitle_Text.font.name = 'Century Gothic'
MockupTitle_Text.font.underline = False
MockupTitle_Text.alignment = PP_ALIGN.CENTER

# WORDING
WORDING = slide.shapes
shape = WORDING.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Cm(18.91), Cm(7.2), Cm(13.83), Cm(1.24)
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(166, 166, 166)

line = shape.line
line.color.rgb = RGBColor(166, 166, 166)
line.width = Pt(2.0)

# WORDINGTitle_Text
WORDINGTitle = slide.shapes.add_textbox(Cm(18.91), Cm(5.11), Cm(13.83), Cm(1.24))
WORDINGTitle_Text = WORDINGTitle.text_frame.add_paragraph()
WORDINGTitle_Text.text = 'WORDING'
WORDINGTitle_Text.font.bold = False
WORDINGTitle_Text.font.italic = False
WORDINGTitle_Text.font.size = Pt(35.4)
WORDINGTitle_Text.font.color.rgb = RGBColor(255, 255, 255)
WORDINGTitle_Text.font.name = 'Century Gothic'
WORDINGTitle_Text.font.underline = False
WORDINGTitle_Text.alignment = PP_ALIGN.CENTER

# WORDINGText

WORDINGText = slide.shapes.add_textbox(Cm(18.91), Cm(6.94), Cm(13.8), Cm(1.1))
WORDINGText_Text = WORDINGText.text_frame.add_paragraph()
WORDINGText_Text.text = Comment
WORDINGText_Text.font.size = Pt(28)
WORDINGText_Text.font.color.rgb = RGBColor(0, 0, 0)
WORDINGText_Text.font.name = 'Century Gothic'
WORDINGText_Text.alignment = PP_ALIGN.LEFT
# WORDINGText.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

# LANDINGText


LANDINGText = slide.shapes.add_textbox(Cm(18.4), Cm(29.56), Cm(14.85), Cm(5.04))
LANDINGText_Text = LANDINGText.text_frame.add_paragraph()
LANDINGText_Text.text = URL
LANDINGText_Text.font.size = Pt(28)
LANDINGText_Text.font.color.rgb = RGBColor(0, 0, 255)
LANDINGText_Text.font.name = 'Century Gothic'
LANDINGText_Text.alignment = PP_ALIGN.CENTER

# LANDING
LANDING = slide.shapes
shape = LANDING.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Cm(18.91), Cm(29.15), Cm(13.83), Cm(1.24)
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(166, 166, 166)

line = shape.line
line.color.rgb = RGBColor(166, 166, 166)
line.width = Pt(2.0)

# LANDINGTitle_Text
LANDINGTitle = slide.shapes.add_textbox(Cm(18.91), Cm(27.06), Cm(13.83), Cm(1.24))
LANDINGTitle_Text = LANDINGTitle.text_frame.add_paragraph()
LANDINGTitle_Text.text = 'LANDING PAGE'
LANDINGTitle_Text.font.bold = False
LANDINGTitle_Text.font.italic = False
LANDINGTitle_Text.font.size = Pt(36)
LANDINGTitle_Text.font.color.rgb = RGBColor(255, 255, 255)
LANDINGTitle_Text.font.name = 'Century Gothic'
LANDINGTitle_Text.font.underline = False
LANDINGTitle_Text.alignment = PP_ALIGN.CENTER

LANDINGTitle_Text2 = LANDINGTitle.text_frame.add_paragraph()
LANDINGTitle_Text2.text = '▼'
LANDINGTitle_Text2.font.size = Pt(16)
LANDINGTitle_Text2.font.color.rgb = RGBColor(127, 127, 127)
LANDINGTitle_Text2.font.name = 'Century Gothic'
LANDINGTitle_Text2.alignment = PP_ALIGN.CENTER

# RightCenter_Text
RightCenter = slide.shapes.add_textbox(Cm(2.31), Cm(18.88), Cm(14.59), Cm(3.89))
RightCenter_Text = RightCenter.text_frame.add_paragraph()
RightCenter_Text.text = "It's the real preview of the visual in the \nWeibo Newsfeed."
RightCenter_Text.font.size = Pt(20)
RightCenter_Text.font.color.rgb = RGBColor(255, 0, 0)
RightCenter_Text.font.name = 'Century Gothic'
RightCenter_Text.alignment = PP_ALIGN.LEFT

# Original Visual Below
OriginalVisualBelow = slide.shapes.add_textbox(Cm(2.71), Cm(25.96), Cm(13.8), Cm(1.1))
OriginalVisualBelow_Text = OriginalVisualBelow.text_frame.add_paragraph()
OriginalVisualBelow_Text.text = 'Original visual below:'
OriginalVisualBelow_Text.font.size = Pt(20)
OriginalVisualBelow_Text.font.color.rgb = RGBColor(0, 0, 0)
OriginalVisualBelow_Text.font.name = 'Century Gothic'
OriginalVisualBelow_Text.alignment = PP_ALIGN.LEFT

# DIGITAL MEDIA TEAM
DIGITALMEDIATEAM = slide.shapes.add_textbox(Cm(0.55), Cm(36.6), Cm(13.8), Cm(1.1))
DIGITALMEDIATEAM_Text = DIGITALMEDIATEAM.text_frame.add_paragraph()
DIGITALMEDIATEAM_Text.text = 'DIGITAL MEDIA TEAM'
DIGITALMEDIATEAM_Text.font.size = Pt(19.7)
DIGITALMEDIATEAM_Text.font.color.rgb = RGBColor(0, 0, 0)
DIGITALMEDIATEAM_Text.font.name = 'Century Gothic'
DIGITALMEDIATEAM_Text.alignment = PP_ALIGN.LEFT

# CenterLine
CenterLine = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER, Cm(34.4), Cm(0.2), Cm(0), Cm(33.4))
CLfill = CenterLine.fill
CLfill.solid()
CLfill.background()
CLline = CenterLine.line
CLline.color.rgb = RGBColor(191, 191, 191)
CLline.width = Pt(0.3)

# LOCAL PUBLICATIONS
LOCALPUBLICATIONS = slide.shapes.add_textbox(Cm(35.74), Cm(7.8), Cm(13.47), Cm(1.6))
LOCALPUBLICATIONS_Text = LOCALPUBLICATIONS.text_frame.add_paragraph()
LOCALPUBLICATIONS_Text.text = 'LOCAL PUBLICATIONS'
LOCALPUBLICATIONS_Text.font.size = Pt(31.5)
LOCALPUBLICATIONS_Text.font.color.rgb = RGBColor(0, 0, 0)
LOCALPUBLICATIONS_Text.font.name = 'Century Gothic'
LOCALPUBLICATIONS_Text.font.underline = True
LOCALPUBLICATIONS_Text.alignment = PP_ALIGN.LEFT

# P1 MD
P1MD = slide.shapes.add_textbox(Cm(41.13), Cm(11.28), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[0]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P1 Comment
P1C = slide.shapes.add_textbox(Cm(40.66), Cm(16.56), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[0]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P2 MD
P1MD = slide.shapes.add_textbox(Cm(51.23), Cm(11.26), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[1]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P2 Comment
P1C = slide.shapes.add_textbox(Cm(50.94), Cm(16.56), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[1]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P3 MD
P1MD = slide.shapes.add_textbox(Cm(61.32), Cm(11.26), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[2]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P3 Comment
P1C = slide.shapes.add_textbox(Cm(60.65), Cm(16.56), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[2]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P4 MD
P1MD = slide.shapes.add_textbox(Cm(41.13), Cm(17.98), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[3]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P4 Comment
P1C = slide.shapes.add_textbox(Cm(40.66), Cm(23.21), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[3]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P5 MD
P1MD = slide.shapes.add_textbox(Cm(51.23), Cm(17.98), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[4]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P5 Comment
P1C = slide.shapes.add_textbox(Cm(50.94), Cm(23.21), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[4]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P6 MD
P1MD = slide.shapes.add_textbox(Cm(61.32), Cm(17.98), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[5]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P6 Comment
P1C = slide.shapes.add_textbox(Cm(60.65), Cm(23.21), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[5]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P7 MD
P1MD = slide.shapes.add_textbox(Cm(41.13), Cm(24.35), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[6]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P7 Comment
P1C = slide.shapes.add_textbox(Cm(41.12), Cm(29.47), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[6]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P8 MD
P1MD = slide.shapes.add_textbox(Cm(51.23), Cm(24.35), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[7]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P8 Comment
P1C = slide.shapes.add_textbox(Cm(51.45), Cm(29.47), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[7]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT
# P9 MD
P1MD = slide.shapes.add_textbox(Cm(61.32), Cm(24.35), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[8]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P9 Comment
P1C = slide.shapes.add_textbox(Cm(61.48), Cm(29.47), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[8]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = "Century Gothic"
P1C_Text.alignment = PP_ALIGN.LEFT
# P10 MD
P1MD = slide.shapes.add_textbox(Cm(41.13), Cm(30.59), Cm(4), Cm(1.18))
P1MD_Text = P1MD.text_frame.add_paragraph()
P1MD_Text.text = MdList[9]
P1MD_Text.font.size = Pt(21.7)
P1MD_Text.font.name = 'Century Gothic'
P1MD_Text.alignment = PP_ALIGN.LEFT
# P10 Comment
P1C = slide.shapes.add_textbox(Cm(40.86), Cm(34.83), Cm(5.6), Cm(3.3))
P1C_Text = P1C.text_frame.add_paragraph()
P1C_Text.text = PostList[9]
P1C_Text.font.size = Pt(16)
P1C_Text.font.name = 'Century Gothic'
P1C_Text.alignment = PP_ALIGN.LEFT

prs.save('test1.pptx')
end_time = time()
run_time = end_time - begin_time
print('PPT生成共用时：', run_time)
